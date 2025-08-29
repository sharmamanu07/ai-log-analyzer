#!/usr/bin/env python3
"""
PowerPoint Report Generator for Security Log Analysis
Creates professional security assessment presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import io
import base64
from datetime import datetime
from typing import List, Dict
import pandas as pd

class SecurityReportGenerator:
    """Generate comprehensive security assessment PowerPoint reports"""
    
    def __init__(self):
        self.prs = Presentation()
        self.setup_theme()
    
    def setup_theme(self):
        """Setup presentation theme and colors"""
        self.colors = {
            'primary': RGBColor(0, 51, 102),      # Dark blue
            'secondary': RGBColor(51, 102, 153),   # Medium blue
            'accent': RGBColor(255, 102, 0),       # Orange
            'success': RGBColor(46, 125, 50),      # Green
            'warning': RGBColor(255, 152, 0),      # Amber
            'danger': RGBColor(198, 40, 40),       # Red
            'critical': RGBColor(136, 14, 79),     # Purple
            'text': RGBColor(33, 33, 33),          # Dark gray
            'background': RGBColor(248, 249, 250)  # Light gray
        }
    
    def generate_report(self, summary: Dict, anomalies: List, logs: List) -> bytes:
        """Generate complete security report"""
        
        # Create slides
        self.create_title_slide(summary)
        self.create_executive_summary_slide(summary, anomalies)
        self.create_risk_assessment_slide(summary, anomalies)
        self.create_threat_overview_slide(anomalies)
        self.create_detailed_findings_slides(anomalies)
        self.create_timeline_slide(anomalies)
        self.create_recommendations_slide(summary)
        self.create_appendix_slide(logs, anomalies)
        # Added addiitonal Slides
        self.create_severity_chart_slide(summary)
        self.create_anomaly_trend_slide(anomalies)
        
        # Save to bytes
        ppt_buffer = io.BytesIO()
        self.prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer.getvalue()
    
    def create_title_slide(self, summary: Dict):
        """Create professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        # Title
        title = slide.shapes.title
        title.text = "Security Log Analysis Report"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title.text_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        subtitle = slide.placeholders[1]
        current_time = datetime.now().strftime("%B %d, %Y")
        risk_level = summary.get('overall_risk', 'UNKNOWN')
        
        subtitle.text = f"Threat Assessment & Security Analysis\n{current_time}\nRisk Level: {risk_level}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']
        
        # Add security classification footer
        left = Inches(0.5)
        top = Inches(7)
        width = Inches(9)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "CONFIDENTIAL - Internal Security Assessment"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = self.colors['text']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def create_executive_summary_slide(self, summary: Dict, anomalies: List):
        """Create executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        self.format_title(title)
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        stats = summary.get('statistics', {})
        
        # Key metrics
        tf.text = f"• Analysis Period: Last 7 days"
        
        p = tf.add_paragraph()
        p.text = f"• Total Anomalies Detected: {stats.get('total_anomalies', 0)}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"• Critical/High Severity Issues: {stats.get('high_severity', 0)}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"• Overall Risk Assessment: {summary.get('overall_risk', 'UNKNOWN')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "• Primary Concerns:"
        p.level = 0
        
        # Top 3 anomaly types
        if anomalies:
            anomaly_types = {}
            for anomaly in anomalies:
                anomaly_types[anomaly.type] = anomaly_types.get(anomaly.type, 0) + 1
            
            sorted_types = sorted(anomaly_types.items(), key=lambda x: x[1], reverse=True)[:3]
            
            for anomaly_type, count in sorted_types:
                p = tf.add_paragraph()
                p.text = f"○ {anomaly_type} ({count} instances)"
                p.level = 1
        
        # Format content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.colors['text']
    
    def create_risk_assessment_slide(self, summary: Dict, anomalies: List):
        """Create risk assessment slide with visual indicators"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_shape.text_frame
        title_tf.text = "Risk Assessment Matrix"
        self.format_title_text(title_tf.paragraphs[0])
        
        # Risk level indicator
        risk_level = summary.get('overall_risk', 'UNKNOWN')
        risk_colors = {
            'LOW': self.colors['success'],
            'MEDIUM': self.colors['warning'],
            'HIGH': self.colors['danger'],
            'CRITICAL': self.colors['critical']
        }
        
        # Risk level box
        left = Inches(1)
        top = Inches(2)
        width = Inches(7)
        height = Inches(1.5)
        
        risk_shape = slide.shapes.add_textbox(left, top, width, height)
        risk_tf = risk_shape.text_frame
        risk_tf.text = f"CURRENT RISK LEVEL: {risk_level}"
        risk_tf.paragraphs[0].font.size = Pt(32)
        risk_tf.paragraphs[0].font.bold = True
        risk_tf.paragraphs[0].font.color.rgb = risk_colors.get(risk_level, self.colors['text'])
        risk_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Severity breakdown
        stats = summary.get('statistics', {})
        severity_data = [
            ('Critical/High', stats.get('high_severity', 0), self.colors['danger']),
            ('Medium', stats.get('medium_severity', 0), self.colors['warning']),
            ('Low', stats.get('low_severity', 0), self.colors['success'])
        ]
        
        y_pos = 4
        for severity, count, color in severity_data:
            left = Inches(2)
            top = Inches(y_pos)
            width = Inches(5)
            height = Inches(0.8)
            
            sev_shape = slide.shapes.add_textbox(left, top, width, height)
            sev_tf = sev_shape.text_frame
            sev_tf.text = f"{severity}: {count} issues"
            sev_tf.paragraphs[0].font.size = Pt(18)
            sev_tf.paragraphs[0].font.color.rgb = color
            
            y_pos += 0.9
    
    def create_threat_overview_slide(self, anomalies: List):
        """Create threat overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Threat Landscape Overview"
        self.format_title(title)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        if not anomalies:
            tf.text = "No significant threats detected in the analyzed period."
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = self.colors['success']
            return
        
        # Group anomalies by type and severity
        threat_summary = {}
        for anomaly in anomalies:
            key = (anomaly.type, anomaly.severity)
            if key not in threat_summary:
                threat_summary[key] = []
            threat_summary[key].append(anomaly)
        
        tf.text = "Identified Threat Categories:"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        
        for (threat_type, severity), threat_list in sorted(threat_summary.items()):
            p = tf.add_paragraph()
            severity_icon = self.get_severity_icon(severity)
            p.text = f"{severity_icon} {threat_type} ({len(threat_list)} instances) - {severity}"
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = self.get_severity_color(severity)
            
            # Add sample description
            if threat_list:
                sample_desc = threat_list[0].description[:80] + "..." if len(threat_list[0].description) > 80 else threat_list[0].description
                p2 = tf.add_paragraph()
                p2.text = f"Example: {sample_desc}"
                p2.level = 1
                p2.font.size = Pt(12)
                p2.font.color.rgb = self.colors['text']
    
    def create_detailed_findings_slides(self, anomalies: List):
        """Create detailed findings slides for top anomalies"""
        
        # Sort anomalies by severity and confidence
        severity_order = {'CRITICAL': 4, 'HIGH': 3, 'MEDIUM': 2, 'LOW': 1}
        sorted_anomalies = sorted(
            anomalies, 
            key=lambda x: (severity_order.get(x.severity, 0), x.confidence), 
            reverse=True
        )
        
        # Create slides for top 5 anomalies
        for i, anomaly in enumerate(sorted_anomalies[:5], 1):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
            
            title = slide.shapes.title
            title.text = f"Finding #{i}: {anomaly.type}"
            self.format_title(title)
            
            content = slide.placeholders[1]
            tf = content.text_frame
            
            # Severity and confidence
            tf.text = f"Severity: {anomaly.severity}"
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.get_severity_color(anomaly.severity)
            
            p = tf.add_paragraph()
            p.text = f"Confidence Level: {anomaly.confidence:.1%}"
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = f"Occurrence Count: {anomaly.count}"
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = f"Time Range: {anomaly.first_seen.strftime('%Y-%m-%d %H:%M')} to {anomaly.last_seen.strftime('%Y-%m-%d %H:%M')}"
            p.font.size = Pt(14)
            
            # Description
            p = tf.add_paragraph()
            p.text = "Description:"
            p.font.size = Pt(14)
            p.font.bold = True
            
            p = tf.add_paragraph()
            p.text = anomaly.description
            p.font.size = Pt(12)
            p.level = 1
            
            # Affected resources
            if anomaly.affected_resources:
                p = tf.add_paragraph()
                p.text = "Affected Resources:"
                p.font.size = Pt(14)
                p.font.bold = True
                
                for resource in anomaly.affected_resources[:5]:  # Limit to 5
                    p = tf.add_paragraph()
                    p.text = f"• {resource}"
                    p.font.size = Pt(12)
                    p.level = 1
    
    def create_timeline_slide(self, anomalies: List):
        """Create timeline visualization slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_shape.text_frame
        title_tf.text = "Security Incident Timeline"
        self.format_title_text(title_tf.paragraphs[0])
        
        if not anomalies:
            # No incidents message
            left = Inches(2)
            top = Inches(3)
            width = Inches(5)
            height = Inches(2)
            
            msg_shape = slide.shapes.add_textbox(left, top, width, height)
            msg_tf = msg_shape.text_frame
            msg_tf.text = "No security incidents detected in the analyzed period."
            msg_tf.paragraphs[0].font.size = Pt(18)
            msg_tf.paragraphs[0].font.color.rgb = self.colors['success']
            msg_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            return
        
        # Sort anomalies by time
        sorted_anomalies = sorted(anomalies, key=lambda x: x.first_seen)
        
        # Create simple timeline text
        y_pos = 2
        for anomaly in sorted_anomalies[:8]:  # Limit to 8 for space
            left = Inches(1)
            top = Inches(y_pos)
            width = Inches(8)
            height = Inches(0.6)
            
            timeline_shape = slide.shapes.add_textbox(left, top, width, height)
            timeline_tf = timeline_shape.text_frame
            
            time_str = anomaly.first_seen.strftime('%Y-%m-%d %H:%M')
            timeline_tf.text = f"{time_str} | {anomaly.type} ({anomaly.severity}) - {anomaly.count} occurrences"
            timeline_tf.paragraphs[0].font.size = Pt(12)
            timeline_tf.paragraphs[0].font.color.rgb = self.get_severity_color(anomaly.severity)
            
            y_pos += 0.7

    # Additional SLide Logic for Anomaly Slide and Severity Chart

    def create_severity_chart_slide(self, summary: Dict):
        """Add pie chart for severity distribution"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank slide

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_tf = title_shape.text_frame
        title_tf.text = "Anomalies by Severity"
        self.format_title_text(title_tf.paragraphs[0])

        stats = summary.get('statistics', {})
        chart_data = CategoryChartData()
        chart_data.categories = ["Critical/High", "Medium", "Low"]
        chart_data.add_series("Anomalies", (
            stats.get('high_severity', 0),
            stats.get('medium_severity', 0),
            stats.get('low_severity', 0)
        ))

        x, y, cx, cy = Inches(2), Inches(1.5), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
        chart.has_legend = True
        chart.legend.include_in_layout = False

    def create_anomaly_trend_slide(self, anomalies: List):
        """Line chart of anomalies over time"""
        if not anomalies:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_tf = title_shape.text_frame
        title_tf.text = "Anomaly Trend Over Time"
        self.format_title_text(title_tf.paragraphs[0])

        # Group anomalies per day
        df = pd.DataFrame([a.first_seen for a in anomalies], columns=["date"])
        df["date"] = df["date"].dt.date
        counts = df["date"].value_counts().sort_index()

        chart_data = CategoryChartData()
        chart_data.categories = list(counts.index)
        chart_data.add_series("Anomalies", list(counts.values))

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
        chart.has_legend = False
 
    
    def create_recommendations_slide(self, summary: Dict):
        """Create recommendations slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Security Recommendations"
        self.format_title(title)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        recommendations = summary.get('recommendations', [])
        
        if not recommendations:
            tf.text = "No specific recommendations at this time. Continue monitoring system logs regularly."
            tf.paragraphs[0].font.size = Pt(16)
            return
        
        tf.text = "Immediate Action Items:"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.colors['primary']
        
        for i, recommendation in enumerate(recommendations[:8], 1):  # Limit to 8
            p = tf.add_paragraph()
            p.text = f"{i}. {recommendation}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text']
            p.level = 0
        
        # Add general recommendations
        p = tf.add_paragraph()
        p.text = "General Security Best Practices:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        general_recs = [
            "Implement continuous log monitoring",
            "Regular security assessments and penetration testing",
            "Keep systems and software updated",
            "Employee security awareness training",
            "Incident response plan testing"
        ]
        
        for rec in general_recs:
            p = tf.add_paragraph()
            p.text = f"• {rec}"
            p.font.size = Pt(12)
            p.level = 1
    
    def create_appendix_slide(self, logs: List, anomalies: List):
        """Create technical appendix slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Technical Appendix"
        self.format_title(title)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Analysis statistics
        tf.text = "Analysis Details:"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"• Total log entries analyzed: {len(logs):,}"
        p.font.size = Pt(14)
        
        if logs:
            time_range = max(log.timestamp for log in logs) - min(log.timestamp for log in logs)
            p = tf.add_paragraph()
            p.text = f"• Analysis time span: {time_range.days} days"
            p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = f"• Anomaly detection algorithms used: Statistical analysis, Pattern matching, Machine learning clustering"
        p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = f"• Report generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S UTC')}"
        p.font.size = Pt(14)
        
        # Detection methodology
        p = tf.add_paragraph()
        p.text = "Detection Methodology:"
        p.font.size = Pt(16)
        p.font.bold = True
        
        methods = [
            "Failed authentication attempt clustering",
            "Error rate spike detection",
            "Unusual access pattern analysis", 
            "Time-based anomaly detection",
            "IP reputation and behavioral analysis"
        ]
        
        for method in methods:
            p = tf.add_paragraph()
            p.text = f"• {method}"
            p.font.size = Pt(12)
            p.level = 1
    
    # Helper methods
    def format_title(self, title_shape):
        """Format title with consistent styling"""
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title_shape.text_frame.paragraphs[0].font.bold = True
    
    def format_title_text(self, paragraph):
        """Format title text paragraph"""
        paragraph.font.size = Pt(32)
        paragraph.font.color.rgb = self.colors['primary']
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
    
    def get_severity_color(self, severity: str):
        """Get color for severity level"""
        colors = {
            'CRITICAL': self.colors['critical'],
            'HIGH': self.colors['danger'],
            'MEDIUM': self.colors['warning'],
            'LOW': self.colors['success']
        }
        return colors.get(severity, self.colors['text'])
    
    def get_severity_icon(self, severity: str):
        """Get icon for severity level"""
        icons = {
            'CRITICAL': '🔴',
            'HIGH': '🟠', 
            'MEDIUM': '🟡',
            'LOW': '🟢'
        }
        return icons.get(severity, '⚪')

def create_security_report(summary: Dict, anomalies: List, logs: List) -> bytes:
    """Create and return PowerPoint security report"""
    generator = SecurityReportGenerator()
    return generator.generate_report(summary, anomalies, logs)

# Example usage
if __name__ == "__main__":
    # Sample data for testing
    sample_summary = {
        'overall_risk': 'HIGH',
        'summary': 'Multiple security anomalies detected requiring immediate attention.',
        'recommendations': [
            'Implement account lockout policies',
            'Review firewall configurations', 
            'Increase monitoring frequency'
        ],
        'statistics': {
            'total_anomalies': 15,
            'high_severity': 5,
            'medium_severity': 7,
            'low_severity': 3
        }
    }
    
    print("PowerPoint report generator ready.")
